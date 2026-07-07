"""
PowerPoint export for the Portfolio tab — "Print Memo Report" button.

Builds a 3-slide deck that mirrors the ATL_Health template structure: each
slide carries one or two PNG screenshots captured from the live browser
view of the Portfolio tab, plus the Xponance branding (gradient bar at the
top, logo at the bottom-left, page number at the bottom-right).

By embedding screenshots rather than rebuilding the panels as native shapes,
the deck looks pixel-identical to what the user sees on screen — which is
exactly what the reference deck does. The trade-off is that the embedded
content isn't editable in PowerPoint; if the user wants to change a number,
they edit it in the tool and re-export.

Slide layout (positions taken directly from ATL_Health.pptx):

    Slide 1
      ├── Portfolio Managers table image  → at (0.45", 1.13"), 9.11" x 2.96"
      └── V-G Positioning image            → at (0.45", 4.08"), 9.09" x 2.62"
    Slide 2
      ├── FactSet Risk Exposures image    → at (1.27", 1.14"), 7.46" x 2.85"
      └── MCR bar chart image              → at (0.00", 3.93"), 10.00" x 2.85"
    Slide 3
      └── Market Cycle (current+proposed)  → at (0.00", 1.67"), 10.00" x 4.46"

Input payload:
  {
    "client_name":    str,
    "benchmark_name": str | null,
    "images": {
      "portfolio_managers": "data:image/png;base64,..."  | None,
      "vg_positioning":     "data:image/png;base64,..."  | None,
      "factset_risk":       "data:image/png;base64,..."  | None,
      "mcr":                "data:image/png;base64,..."  | None,
      "market_cycle":       "data:image/png;base64,..."  | None,
    }
  }

If an image is None, the slide that depends on it is omitted from the deck.
The caller is informed via the X-Skipped-Slides header.
"""

from io import BytesIO
import base64
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ── Slide geometry ────────────────────────────────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# Brand colours sampled from the reference deck
COLOR_TEXT   = RGBColor(0x1E, 0x2D, 0x3D)
COLOR_BORDER = RGBColor(0xE2, 0xE7, 0xED)

# Logo asset path (X-glyph PNG, with "ponance" rendered next to it as text)
ASSETS_DIR = os.path.join(os.path.dirname(__file__), 'static', 'assets')
LOGO_PATH  = os.path.join(ASSETS_DIR, 'xponance_logo.png')


# ── Helpers ───────────────────────────────────────────────────────────────
def _decode_data_url(data_url):
    """Strip the 'data:image/png;base64,' prefix and decode. Returns bytes
    suitable to feed to add_picture, or None if the input isn't a valid
    PNG data URL."""
    if not data_url or not isinstance(data_url, str):
        return None
    m = re.match(r'^data:image/(?:png|jpeg);base64,(.+)$', data_url, re.DOTALL)
    if not m:
        return None
    try:
        return base64.b64decode(m.group(1))
    except Exception:
        return None


def _add_branding(slide, slide_no):
    """Draw the chrome that goes on every content slide:
       - thin horizontal gradient bar at top
       - Xponance logo (X glyph PNG + 'ponance' text) at bottom-left
       - slide number at bottom-right

    The gradient bar XML must be inserted in OOXML schema order
    (xfrm → prstGeom → fill → ln → effects); inserting it after the
    existing <a:ln> element triggers PowerPoint's repair dialog.
    """
    # ── Top gradient bar ──────────────────────────────────────────────────
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.0), Inches(0.30),
                                  SLIDE_W, Inches(0.04))
    bar.line.fill.background()
    sp_pr = bar._element.spPr
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for tag in ('noFill', 'solidFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill'):
        for el in sp_pr.findall(qn(f'a:{tag}')):
            sp_pr.remove(el)
    grad_xml = (
        f'<a:gradFill rotWithShape="1" xmlns:a="{a_ns}">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="0C4BA3"/></a:gs>'
        '<a:gs pos="50000"><a:srgbClr val="1E88E5"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="4FC3F7"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    grad_el = etree.fromstring(grad_xml)
    ln_el = sp_pr.find(qn('a:ln'))
    if ln_el is not None:
        ln_el.addprevious(grad_el)
    else:
        sp_pr.append(grad_el)

    # ── Bottom-left logo: full "Xponance" wordmark as a single image.
    # The image (static/assets/xponance_logo.png) is the brand logo file the
    # user supplied — keep its native aspect ratio (~2.84:1) so it doesn't
    # squash. At 1.30" wide, height is 1.30/2.84 ≈ 0.46".
    logo_w = Inches(1.30)
    logo_h = Inches(0.46)
    logo_left = Inches(0.30)
    logo_top  = SLIDE_H - logo_h - Inches(0.18)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, logo_left, logo_top, logo_w, logo_h)

    # ── Bottom-right slide number ─────────────────────────────────────────
    num_box = slide.shapes.add_textbox(
        SLIDE_W - Inches(0.50), SLIDE_H - Inches(0.35),
        Inches(0.30), Inches(0.25))
    tf = num_box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(slide_no)
    run.font.name = 'Calibri'
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _embed_image(slide, image_bytes, left, top, width, height):
    """Embed a PNG at the given position. python-pptx requires a file-like
    object, so we wrap the bytes in BytesIO."""
    if not image_bytes:
        return None
    return slide.shapes.add_picture(BytesIO(image_bytes),
                                     left, top, width, height)


# ── Public API ────────────────────────────────────────────────────────────
def build_portfolio_pptx(payload):
    """Build the 3-slide memo deck from the captured browser screenshots.
    Returns (bytes, summary_dict). summary_dict.skipped lists any sections
    that were omitted because their image wasn't provided."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    images = payload.get('images') or {}
    img_managers  = _decode_data_url(images.get('portfolio_managers'))
    img_vg        = _decode_data_url(images.get('vg_positioning'))
    img_factset   = _decode_data_url(images.get('factset_risk'))
    img_mcr       = _decode_data_url(images.get('mcr'))
    img_market    = _decode_data_url(images.get('market_cycle'))

    skipped = []
    slide_no = 0

    # ── Slide 1: Portfolio Managers + V-G ─────────────────────────────────
    if img_managers or img_vg:
        slide_no += 1
        slide = _blank_slide(prs)
        _add_branding(slide, slide_no)
        if img_managers:
            _embed_image(slide, img_managers,
                         Inches(0.45), Inches(1.13),
                         Inches(9.11), Inches(2.96))
        else:
            skipped.append('Slide 1: portfolio managers panel not captured')
        if img_vg:
            _embed_image(slide, img_vg,
                         Inches(0.45), Inches(4.08),
                         Inches(9.09), Inches(2.62))
        else:
            skipped.append('Slide 1: V-G positioning panel not captured')
    else:
        skipped.append('Slide 1: skipped - no managers or V-G image')

    # ── Slide 2: FactSet Risk + MCR ───────────────────────────────────────
    if img_factset or img_mcr:
        slide_no += 1
        slide = _blank_slide(prs)
        _add_branding(slide, slide_no)
        if img_factset:
            _embed_image(slide, img_factset,
                         Inches(1.27), Inches(1.14),
                         Inches(7.46), Inches(2.85))
        else:
            skipped.append('Slide 2: FactSet risk panel not captured - load risk file')
        if img_mcr:
            _embed_image(slide, img_mcr,
                         Inches(0.00), Inches(3.93),
                         Inches(10.00), Inches(2.85))
        else:
            skipped.append('Slide 2: MCR chart not captured - load factor returns')
    else:
        skipped.append('Slide 2: skipped - no FactSet or MCR image')

    # ── Slide 3: Market Cycle ─────────────────────────────────────────────
    if img_market:
        slide_no += 1
        slide = _blank_slide(prs)
        _add_branding(slide, slide_no)
        _embed_image(slide, img_market,
                     Inches(0.00), Inches(1.67),
                     Inches(10.00), Inches(4.46))
    else:
        skipped.append('Slide 3: market cycle not captured - load universe file')

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue(), {'skipped': skipped, 'n_slides': len(prs.slides)}
